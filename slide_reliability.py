from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# =============================================================================
# Helper function to add a title and content slide
# =============================================================================
def add_title_and_content_slide(prs, title_text, content_items):
    """
    Adds a new slide with a title and bulleted content.
    """
    slide_layout = prs.slide_layouts[1]  # Layout 'Title and Content'
    slide = prs.slides.add_slide(slide_layout)
    
    # Set Title
    title = slide.shapes.title
    title.text = title_text
    
    # Set Content
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()  # Clear existing content

    for item in content_items:
        if isinstance(item, list): # Handle nested bullets
            p = tf.add_paragraph()
            p.text = item[0]
            p.level = 1
            for sub_item in item[1:]:
                p_sub = tf.add_paragraph()
                p_sub.text = sub_item
                p_sub.level = 2
        else:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            
# =============================================================================
# Main function to generate the presentation
# =============================================================================
def main():
    """
    Generates the Service Reliability PowerPoint presentation.
    """
    # Create a presentation object
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]  # Layout 'Title Slide'
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Service Reliability"
    subtitle.text = "Xây dựng và Đo lường Dịch vụ Đáng Tin cậy"

    # --- Slide 2: Agenda ---
    add_title_and_content_slide(
        prs,
        "Nội dung chính",
        [
            "Service Reliability là gì?",
            "Tại sao Độ tin cậy lại quan trọng?",
            "Các Metric chính để đo lường (The Four Golden Signals)",
            "SLI, SLO, và SLA: Nền tảng của Reliability",
            "Các Metric bổ trợ: MTTR & MTBF",
            "Tổng kết"
        ]
    )

    # --- Slide 3: What is Service Reliability? ---
    add_title_and_content_slide(
        prs,
        "Service Reliability là gì?",
        [
            "Là khả năng một dịch vụ hoạt động ĐÚNG CHỨC NĂNG, NHẤT QUÁN và LIÊN TỤC.",
            "Đây là lời hứa với người dùng: \"Dịch vụ sẽ hoạt động như bạn mong đợi, bất cứ khi nào bạn cần.\"",
            "Khác biệt với Availability (Độ sẵn sàng):",
            [
                "Availability: Dịch vụ có 'up' không?",
                "Reliability: Dịch vụ có 'up' VÀ hoạt động đúng, đúng hiệu suất không?"
            ],
            "Ví dụ: API luôn phản hồi nhưng không xử lý giao dịch => Available nhưng không Reliable."
        ]
    )

    # --- Slide 4: Why is it Important? ---
    add_title_and_content_slide(
        prs,
        "Tại sao Độ tin cậy lại quan trọng?",
        [
            "Xây dựng lòng tin của khách hàng: Người dùng sẽ rời bỏ các dịch vụ không ổn định.",
            "Tác động trực tiếp đến doanh thu: Downtime = Mất tiền.",
            "Nâng cao năng suất cho đội ngũ kỹ sư: Giảm thời gian chữa cháy (firefighting), tập trung vào phát triển tính năng mới.",
            "Bảo vệ danh tiếng: Một sự cố lớn có thể gây tổn hại hình ảnh của team."
        ]
    )

    # --- Slide 5: The Four Golden Signals ---
    add_title_and_content_slide(
        prs,
        "The Four Golden Signals",
        [
            "Bộ tứ metric quan trọng nhất để giám sát sức khỏe của một hệ thống, theo Google SRE.",
            "1. Latency (Độ trễ)",
            "2. Traffic (Lưu lượng)",
            "3. Errors (Tỷ lệ lỗi)",
            "4. Saturation (Độ bão hòa)"
        ]
    )

    # --- Slide 6: Golden Signal 1: Latency ---
    add_title_and_content_slide(
        prs,
        "Golden Signal 1: Latency (Độ trễ)",
        [
            "Đo lường thời gian xử lý một yêu cầu.",
            "\"Chậm cũng là một dạng lỗi.\"",
            "Cách đo lường hiệu quả:",
            [
                "Không dùng giá trị trung bình (average).",
                "Sử dụng phân vị (Percentiles): p50, p95, p99 để hiểu trải nghiệm của số đông và của những người dùng tệ nhất."
            ],
            "Mục tiêu (SLO ví dụ): 99% yêu cầu phải được xử lý dưới 300ms."
        ]
    )

    # --- Slide 7: Latency p95/p99: Industry Standard ---
    add_title_and_content_slide(
        prs,
        "Latency p95/p99: Tiêu chuẩn ngành",
        [
            "Đối với các dịch vụ User-Facing (web/mobile app), p95 hoặc p99 latency là thước đo quan trọng nhất.",
            "Ngưỡng tâm lý người dùng:",
            [
                "< 100ms: Phản hồi tức thì.",
                "200-250ms: Vẫn cảm thấy nhanh, chấp nhận được.",
                "> 500ms: Bắt đầu cảm thấy chậm, khó chịu.",
                "> 1s: Mất tập trung, có nguy cơ rời bỏ."
            ],
            "Vì vậy, 200-250ms cho p95/p99 là tiêu chuẩn ngành cho các tác vụ chính."
        ]
    )

    # --- Slide 8: Golden Signal 2: Traffic ---
    add_title_and_content_slide(
        prs,
        "Golden Signal 2: Traffic (Lưu lượng)",
        [
            "Đo lường mức độ nhu cầu (demand) đang đặt lên hệ thống.",
            "Giúp hiểu rõ quy mô và sự tăng trưởng của dịch vụ.",
            "Cách đo lường:",
            [
                "Đối với API/web service: Requests Per Second (RPS).",
                "Đối với hệ thống streaming: Events Per Second.",
                "Đối với hệ thống e-commerce: Giao dịch mỗi phút."
            ],
            "Quan trọng cho việc capacity planning và phát hiện các thay đổi bất thường (vd: botnet attack)."
        ]
    )

    # --- Slide 9: Golden Signal 3: Errors ---
    add_title_and_content_slide(
        prs,
        "Golden Signal 3: Errors (Tỷ lệ lỗi)",
        [
            "Đo lường tần suất các yêu cầu bị lỗi.",
            "Bao gồm cả lỗi rõ ràng (explicit) và lỗi ngầm (implicit):",
            [
                "Lỗi rõ ràng: HTTP 500, crash...",
                "Lỗi ngầm: Trả về kết quả 200 OK nhưng nội dung sai."
            ],
            "Công thức: (Số request lỗi / Tổng số request) x 100%",
            "Mục tiêu (SLO ví dụ): Tỷ lệ lỗi cho API thanh toán phải < 0.05%."
        ]
    )

    # --- Slide 10: Golden Signal 4: Saturation ---
    add_title_and_content_slide(
        prs,
        "Golden Signal 4: Saturation (Độ bão hòa)",
        [
            "Đo lường mức độ \"bận rộn\" hoặc \"đầy\" của hệ thống.",
            "Là một chỉ số dự báo (leading indicator) về các sự cố sắp xảy ra.",
            "Các tài nguyên cần theo dõi:",
            [
                "CPU utilization",
                "Memory usage",
                "Disk I/O",
                "Network bandwidth"
            ],
            "Mục tiêu (SLO ví dụ): Mức sử dụng CPU không được vượt quá 85% trong 5 phút."
        ]
    )

    # --- Slide 11: SLI, SLO, SLA ---
    add_title_and_content_slide(
        prs,
        "SLI, SLO, SLA: Nền tảng của Reliability",
        [
            "SLI (Service Level Indicator): Chỉ số bạn đo lường.",
            [
                "Là một chỉ số định lượng (quantitative measure) về một khía cạnh của dịch vụ.",
                "Phải đo lường được và phản ánh đúng trải nghiệm của người dùng.",
                "Ví dụ: Tỷ lệ lỗi của request, Độ trễ của response, Tỷ lệ request thành công."
            ],
            "SLO (Service Level Objective): Mục tiêu bạn cam kết đạt được.",
            [
                "Là mục tiêu cụ thể cho một SLI trong một khoảng thời gian.",
                "Là cam kết nội bộ, không phải là hợp đồng với khách hàng.",
                "Ví dụ: 99.9% request phải thành công trong 1 tháng.",
                "Error Budget được suy ra từ SLO: 100% - SLO. Ví dụ: 1 - 99.9% = 0.1% lỗi được phép."
            ],
            "SLA (Service Level Agreement): Hợp đồng với khách hàng.",
            [
                "Là một hợp đồng chính thức với khách hàng, bao gồm các SLO.",
                "Thường có điều khoản phạt (penalties) nếu không đáp ứng được.",
                "SLA thường 'lỏng' hơn SLO để giảm rủi ro kinh doanh.",
                "Ví dụ: Nếu uptime < 99.9% trong một tháng, khách hàng sẽ được giảm giá 10%."
            ],
        ]
    )

    # --- Slide 12: Supporting Metrics ---
    add_title_and_content_slide(
        prs,
        "Các Metric bổ trợ quan trọng",
        [
            "MTTR (Mean Time To Repair/Recovery): Thời gian trung bình để phục hồi sau sự cố.",
            [
                "Đo lường khả năng phản ứng của đội ngũ.",
                "Mục tiêu: Càng thấp càng tốt."
            ],
            "MTBF (Mean Time Between Failures): Thời gian trung bình giữa các sự cố.",
            [
                "Đo lường sự ổn định của hệ thống.",
                "Mục tiêu: Càng cao càng tốt."
            ]
        ]
    )

    # --- Slide 13: Summary ---
    add_title_and_content_slide(
        prs,
        "Tổng kết",
        [
            "Reliability là nền tảng của một dịch vụ thành công, tạo dựng niềm tin cho người dùng.",
            "Bắt đầu đo lường với The Four Golden Signals: Latency, Errors, Traffic, Saturation.",
            "Thiết lập SLI và SLO rõ ràng để định nghĩa \"đủ tốt\" là như thế nào.",
            "Luôn cố gắng cải thiện MTTR (phản ứng nhanh hơn) và MTBF (ít lỗi hơn).",
            "Độ tin cậy là trách nhiệm của cả đội ngũ, không chỉ của riêng ai."
        ]
    )

    # --- Slide 14: Q&A ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Q & A"
    subtitle.text = "Cảm ơn đã lắng nghe!"

    # Save the presentation
    file_path = "Service_Reliability_Presentation.pptx"
    prs.save(file_path)

    print(f"Presentation saved to {file_path}")

if __name__ == "__main__":
    main()